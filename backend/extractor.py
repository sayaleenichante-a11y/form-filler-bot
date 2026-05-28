from pathlib import Path
import json
import csv
import fitz
import docx
import pandas as pd
from pptx import Presentation
from striprtf.striprtf import rtf_to_text
from bs4 import BeautifulSoup
from PIL import Image
import pytesseract


MAX_CHARS = 800000
MAX_PDF_PAGES = 500


def safe_limit(text: str) -> str:
    if not text:
        return ""
    return text[:MAX_CHARS]


def extract_text(file_path):
    file_path = Path(file_path)
    ext = file_path.suffix.lower()

    if ext == ".pdf":
        return safe_limit(extract_pdf(file_path))

    if ext == ".docx":
        return safe_limit(extract_docx(file_path))

    if ext in [".txt", ".md", ".log"]:
        return safe_limit(file_path.read_text(encoding="utf-8", errors="ignore"))

    if ext in [".xlsx", ".xls"]:
        return safe_limit(extract_excel(file_path))

    if ext == ".csv":
        return safe_limit(extract_csv(file_path))

    if ext == ".pptx":
        return safe_limit(extract_pptx(file_path))

    if ext == ".rtf":
        return safe_limit(rtf_to_text(file_path.read_text(encoding="utf-8", errors="ignore")))

    if ext in [".html", ".htm"]:
        return safe_limit(extract_html(file_path))

    if ext == ".json":
        return safe_limit(extract_json(file_path))

    if ext in [".png", ".jpg", ".jpeg", ".webp", ".bmp"]:
        return safe_limit(extract_image(file_path))

    raise ValueError(f"Unsupported file type: {ext}")


def extract_pdf(file_path):
    text_parts = []

    doc = fitz.open(file_path)
    total_pages = min(len(doc), MAX_PDF_PAGES)

    for page_index in range(total_pages):
        page = doc[page_index]
        page_text = page.get_text("text") or ""

        if page_text.strip():
            text_parts.append(page_text)
        else:
            try:
                pix = page.get_pixmap(dpi=160)
                temp_img = file_path.parent / f"temp_ocr_page_{page_index}.png"
                pix.save(temp_img)
                ocr_text = pytesseract.image_to_string(Image.open(temp_img))
                text_parts.append(ocr_text)
                temp_img.unlink(missing_ok=True)
            except Exception:
                pass

    return "\n".join(text_parts)


def extract_docx(file_path):
    document = docx.Document(file_path)

    text_parts = []

    for paragraph in document.paragraphs:
        if paragraph.text.strip():
            text_parts.append(paragraph.text.strip())

    for table in document.tables:
        for row in table.rows:
            row_text = []
            for cell in row.cells:
                value = cell.text.strip()
                if value:
                    row_text.append(value)
            if row_text:
                text_parts.append(" | ".join(row_text))

    return "\n".join(text_parts)


def extract_excel(file_path):
    text_parts = []
    excel = pd.ExcelFile(file_path)

    for sheet in excel.sheet_names:
        df = excel.parse(sheet)
        text_parts.append(f"\nSheet: {sheet}\n")
        text_parts.append(df.to_string(index=False))

    return "\n".join(text_parts)


def extract_csv(file_path):
    text_parts = []

    with open(file_path, newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            text_parts.append(" | ".join(row))

    return "\n".join(text_parts)


def extract_pptx(file_path):
    presentation = Presentation(file_path)
    text_parts = []

    for slide_no, slide in enumerate(presentation.slides, start=1):
        text_parts.append(f"\nSlide {slide_no}\n")

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())

    return "\n".join(text_parts)


def extract_html(file_path):
    html = file_path.read_text(encoding="utf-8", errors="ignore")
    soup = BeautifulSoup(html, "html.parser")
    return soup.get_text(separator="\n")


def extract_json(file_path):
    data = json.loads(file_path.read_text(encoding="utf-8", errors="ignore"))
    return json.dumps(data, indent=2, ensure_ascii=False)


def extract_image(file_path):
    return pytesseract.image_to_string(Image.open(file_path))